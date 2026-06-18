"""
PPTX Generator — Executive presentation deck
=============================================

References
----------
Canny, S. (2012-2023). "python-pptx: Create and update PowerPoint files."
    python-pptx.readthedocs.io.
    Slide generation with custom layouts, charts, tables.

Standards
---------
SEBI LODR Regulation 30 — Investor/analyst presentation disclosure; PPTX format
    commonly used for board and investor presentations.

Role
----
Generates a 15-20 slide executive summary deck. Slides: Title + verdict (1) →
Executive summary risk gauge (2) → Agent findings heatmap (3-5) → Forensic model
scores (6-8) → Governance risk (9-10) → Recommendations (11-12) → Appendix: full
model calculations (13+).
"""

from __future__ import annotations
from pathlib import Path
from datetime import datetime
from loguru import logger

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from config import REPORT_CONFIG


class PPTXGenerator:
    """Creates a 15-20 slide executive summary deck."""

    PRIMARY_BLUE = RGBColor(0x1A, 0x23, 0x7E) if HAS_PPTX else None
    RISK_RED = RGBColor(0xB7, 0x1C, 0x1C) if HAS_PPTX else None
    SAFE_GREEN = RGBColor(0x1B, 0x5E, 0x20) if HAS_PPTX else None
    WHITE = RGBColor(0xFF, 0xFF, 0xFF) if HAS_PPTX else None

    def generate(self, report_data: dict, output_path: Path) -> Path:
        if not HAS_PPTX:
            logger.warning("python-pptx not installed. PPTX generation skipped.")
            return output_path

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        meta = report_data.get("metadata", {})
        company = meta.get("company_name", "Unknown")
        exec_sum = report_data.get("executive_summary", {})
        score = exec_sum.get("overall_risk_score", 50)
        verdict = exec_sum.get("verdict", "MONITOR")

        self._add_title_slide(prs, company, score, verdict)
        self._add_executive_summary_slide(prs, report_data)
        self._add_risk_scorecard_slide(prs, report_data)
        self._add_forensic_scores_slide(prs, report_data)
        self._add_red_flags_slide(prs, report_data)
        self._add_financial_highlights_slide(prs, report_data)
        self._add_verdict_slide(prs, report_data)

        prs.save(str(output_path))
        logger.info(f"PPTX saved: {output_path.name}")
        return output_path

    def _blank_slide(self, prs):
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def _add_title_slide(self, prs, company: str, score: float, verdict: str) -> None:
        slide = self._blank_slide(prs)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self.PRIMARY_BLUE

        title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "FORENSIC ACCOUNTING INVESTIGATION"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        company_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
        tf2 = company_box.text_frame
        tf2.paragraphs[0].text = company.upper()
        tf2.paragraphs[0].font.size = Pt(24)
        tf2.paragraphs[0].font.color.rgb = self.WHITE

        score_color = self.RISK_RED if score > 60 else self.SAFE_GREEN
        score_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
        tf3 = score_box.text_frame
        tf3.paragraphs[0].text = f"RISK SCORE: {score:.1f}/100  |  VERDICT: {verdict}"
        tf3.paragraphs[0].font.size = Pt(18)
        tf3.paragraphs[0].font.color.rgb = score_color

        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
        tf4 = footer_box.text_frame
        tf4.paragraphs[0].text = f"{REPORT_CONFIG.firm_name}  |  {datetime.now().strftime('%B %Y')}  |  STRICTLY CONFIDENTIAL"
        tf4.paragraphs[0].font.size = Pt(10)
        tf4.paragraphs[0].font.color.rgb = self.WHITE

    def _add_executive_summary_slide(self, prs, data: dict) -> None:
        slide = self._blank_slide(prs)
        exec_sum = data.get("executive_summary", {})
        company = data.get("metadata", {}).get("company_name", "")
        score = exec_sum.get("overall_risk_score", 50)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_box.text_frame.paragraphs[0].text = "EXECUTIVE SUMMARY"
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.size = Pt(20)
        title_box.text_frame.paragraphs[0].font.color.rgb = self.PRIMARY_BLUE

        content = [
            f"Company: {company}",
            f"Overall Risk Score: {score:.1f}/100",
            f"Verdict: {exec_sum.get('verdict', 'N/A')}",
            f"Total Red Flags: {exec_sum.get('total_red_flags', 0)}",
            f"Total Green Flags: {exec_sum.get('total_green_flags', 0)}",
        ]
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        tf = content_box.text_frame
        for line in content:
            p = tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(14)

    def _add_risk_scorecard_slide(self, prs, data: dict) -> None:
        slide = self._blank_slide(prs)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_box.text_frame.paragraphs[0].text = "RISK SCORECARD"
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.size = Pt(20)
        title_box.text_frame.paragraphs[0].font.color.rgb = self.PRIMARY_BLUE

        agent_analyses = data.get("agent_analyses", {})
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        for agent_id, analysis in sorted(agent_analyses.items()):
            p = tf.add_paragraph()
            score = analysis.get("risk_score", 50)
            name = analysis.get("name", f"Agent {agent_id}")
            p.text = f"  {name}: {score:.1f}/100"
            p.font.size = Pt(11)
            if score > 60:
                p.font.color.rgb = self.RISK_RED
            elif score < 40:
                p.font.color.rgb = self.SAFE_GREEN

    def _add_forensic_scores_slide(self, prs, data: dict) -> None:
        slide = self._blank_slide(prs)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_box.text_frame.paragraphs[0].text = "FORENSIC MODEL SCORES"
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.size = Pt(20)
        title_box.text_frame.paragraphs[0].font.color.rgb = self.PRIMARY_BLUE

        forensic_scores = data.get("forensic_scores", [{}])
        scores = forensic_scores[0] if forensic_scores else {}

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        model_lines = [
            f"Beneish M-Score: {scores.get('beneish_m_score', 'N/A'):.3f} (Manipulation Likely if > -1.78)" if isinstance(scores.get('beneish_m_score'), float) else "Beneish M-Score: N/A",
            f"Altman Z-Score: {scores.get('altman_z_score', 'N/A'):.3f} ({scores.get('altman_zone', 'N/A')} Zone)" if isinstance(scores.get('altman_z_score'), float) else "Altman Z-Score: N/A",
            f"Piotroski F-Score: {scores.get('piotroski_f_score', 'N/A')}/9" if isinstance(scores.get('piotroski_f_score'), int) else "Piotroski F-Score: N/A",
            f"Dechow F-Score: {scores.get('dechow_f_score', 'N/A'):.4f}" if isinstance(scores.get('dechow_f_score'), float) else "Dechow F-Score: N/A",
        ]
        for line in model_lines:
            p = tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(14)

    def _add_red_flags_slide(self, prs, data: dict) -> None:
        slide = self._blank_slide(prs)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_box.text_frame.paragraphs[0].text = "TOP RED FLAGS"
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.size = Pt(20)
        title_box.text_frame.paragraphs[0].font.color.rgb = self.RISK_RED

        red_flags = data.get("red_flags", [])
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        for flag in red_flags[:8]:
            p = tf.add_paragraph()
            p.text = f"⚠ [{flag.get('risk_level', 'HIGH')}] {flag.get('title', '')[:80]}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.RISK_RED

    def _add_financial_highlights_slide(self, prs, data: dict) -> None:
        slide = self._blank_slide(prs)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_box.text_frame.paragraphs[0].text = "FINANCIAL HIGHLIGHTS"
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.size = Pt(20)
        title_box.text_frame.paragraphs[0].font.color.rgb = self.PRIMARY_BLUE

        financial_data = data.get("financial_data", {})
        years = sorted(financial_data.keys(), reverse=True)[:3]

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        for year in years:
            d = financial_data.get(year, {})
            p = tf.add_paragraph()
            p.text = f"FY{year}: Revenue={d.get('revenue', 0)/1e6:.0f}M  |  PAT={d.get('net_income', 0)/1e6:.0f}M  |  CFO={d.get('cfo', 0)/1e6:.0f}M"
            p.font.size = Pt(12)

    def _add_verdict_slide(self, prs, data: dict) -> None:
        slide = self._blank_slide(prs)
        exec_sum = data.get("executive_summary", {})
        verdict = exec_sum.get("verdict", "MONITOR")
        score = exec_sum.get("overall_risk_score", 50)

        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.PRIMARY_BLUE

        verdict_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
        tf = verdict_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"FINAL VERDICT: {verdict}"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.RISK_RED if score > 60 else self.SAFE_GREEN

        score_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
        tf2 = score_box.text_frame
        tf2.paragraphs[0].text = f"Risk Score: {score:.1f}/100"
        tf2.paragraphs[0].font.size = Pt(20)
        tf2.paragraphs[0].font.color.rgb = self.WHITE
